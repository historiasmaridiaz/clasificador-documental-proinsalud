from __future__ import annotations

import csv
import hashlib
import io
import json
import os
import re
import shutil
import unicodedata
import zipfile
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Any, Iterable


TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".rtf", ".log"}
OFFICE_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xlsm", ".xls", ".pptx", ".odt", ".eml"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}
SUPPORTED_DOCUMENT_EXTENSIONS = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | IMAGE_EXTENSIONS
MAX_DOCUMENT_BYTES = 50 * 1024 * 1024
MAX_ARCHIVE_BYTES = 300 * 1024 * 1024
MAX_ARCHIVE_MEMBERS = 1_000
MAX_EXTRACTED_CHARS = 250_000


@dataclass
class DocumentPayload:
    name: str
    data: bytes
    source_modified_at: str = ""

    @property
    def extension(self) -> str:
        return Path(self.name).suffix.lower()

    @property
    def sha256(self) -> str:
        return hashlib.sha256(self.data).hexdigest()


@dataclass
class ExtractedDocument:
    name: str
    text: str
    sha256: str
    extension: str
    metadata: dict[str, Any] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)

    @property
    def preview(self) -> str:
        return self.text[:3_000]


def _coerce_upload(item: Any) -> DocumentPayload:
    if isinstance(item, DocumentPayload):
        return item
    if isinstance(item, tuple) and len(item) == 2:
        return DocumentPayload(str(item[0]), bytes(item[1]))
    if isinstance(item, dict):
        return DocumentPayload(
            str(item["name"]),
            bytes(item["data"]),
            str(item.get("source_modified_at", "")),
        )
    name = str(getattr(item, "name", "archivo"))
    modified = getattr(item, "last_modified", None) or getattr(item, "lastModified", None) or ""
    if hasattr(modified, "isoformat"):
        modified = modified.isoformat()
    if hasattr(item, "getvalue"):
        return DocumentPayload(name, bytes(item.getvalue()), str(modified))
    return DocumentPayload(name, bytes(item.read()), str(modified))


def _safe_member_name(name: str) -> str:
    normalized = name.replace("\\", "/").lstrip("/")
    return "/".join(p for p in normalized.split("/") if p not in ("", ".", ".."))


def expand_document_uploads(items: Iterable[Any]) -> tuple[list[DocumentPayload], list[str]]:
    documents: list[DocumentPayload] = []
    warnings: list[str] = []
    for item in items:
        payload = _coerce_upload(item)
        if payload.extension != ".zip":
            if payload.extension in SUPPORTED_DOCUMENT_EXTENSIONS:
                if len(payload.data) <= MAX_DOCUMENT_BYTES:
                    documents.append(payload)
                else:
                    warnings.append(f"{payload.name}: supera 50 MB y fue omitido.")
            else:
                warnings.append(f"{payload.name}: formato no compatible.")
            continue

        try:
            with zipfile.ZipFile(io.BytesIO(payload.data)) as archive:
                members = [m for m in archive.infolist() if not m.is_dir()]
                if len(members) > MAX_ARCHIVE_MEMBERS:
                    raise ValueError(f"el ZIP supera {MAX_ARCHIVE_MEMBERS} archivos")
                total = 0
                for member in members:
                    member_name = _safe_member_name(member.filename)
                    ext = Path(member_name).suffix.lower()
                    if ext not in SUPPORTED_DOCUMENT_EXTENSIONS:
                        continue
                    if member.file_size > MAX_DOCUMENT_BYTES:
                        warnings.append(f"{member_name}: supera 50 MB y fue omitido.")
                        continue
                    total += member.file_size
                    if total > MAX_ARCHIVE_BYTES:
                        raise ValueError("el contenido descomprimido supera 300 MB")
                    member_modified = "-".join(
                        [
                            f"{member.date_time[0]:04d}",
                            f"{member.date_time[1]:02d}",
                            f"{member.date_time[2]:02d}",
                        ]
                    ) + f"T{member.date_time[3]:02d}:{member.date_time[4]:02d}:{member.date_time[5]:02d}"
                    documents.append(
                        DocumentPayload(
                            f"{payload.name}/{member_name}",
                            archive.read(member),
                            member_modified,
                        )
                    )
        except Exception as exc:
            warnings.append(f"{payload.name}: no se pudo abrir ({exc}).")
    return documents, warnings


def _decode_text(data: bytes) -> str:
    try:
        from charset_normalizer import from_bytes

        match = from_bytes(data).best()
        if match:
            return str(match)
    except Exception:
        pass
    return data.decode("utf-8", errors="replace")


def _candidate_tesseract_commands() -> list[str]:
    candidates: list[str] = []
    configured = os.environ.get("TESSERACT_CMD", "").strip()
    if configured:
        candidates.append(configured)
    located = shutil.which("tesseract")
    if located:
        candidates.append(located)
    if os.name == "nt":
        roots = [
            os.environ.get("PROGRAMFILES", r"C:\Program Files"),
            os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)"),
            os.environ.get("LOCALAPPDATA", ""),
        ]
        relatives = [
            Path("Tesseract-OCR") / "tesseract.exe",
            Path("Programs") / "Tesseract-OCR" / "tesseract.exe",
        ]
        for root in roots:
            if not root:
                continue
            for relative in relatives:
                candidates.append(str(Path(root) / relative))
    unique: list[str] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = os.path.normcase(os.path.abspath(candidate)) if candidate else ""
        if candidate and key not in seen:
            seen.add(key)
            unique.append(candidate)
    return unique


def tesseract_status() -> dict[str, Any]:
    """Informa si el motor OCR local está disponible y configura pytesseract."""
    for command in _candidate_tesseract_commands():
        if Path(command).is_file() or shutil.which(command):
            try:
                import pytesseract

                pytesseract.pytesseract.tesseract_cmd = command
                version = str(pytesseract.get_tesseract_version())
                return {"available": True, "command": command, "version": version}
            except Exception:
                continue
    try:
        import pytesseract

        version = str(pytesseract.get_tesseract_version())
        return {
            "available": True,
            "command": str(getattr(pytesseract.pytesseract, "tesseract_cmd", "tesseract")),
            "version": version,
        }
    except Exception as exc:
        return {"available": False, "command": "", "version": "", "error": str(exc)}


def _prepare_ocr_image(image: Any) -> Any:
    from PIL import ImageOps

    prepared = ImageOps.exif_transpose(image).convert("RGB")
    if prepared.width < 1_600:
        factor = min(2.5, 1_600 / max(1, prepared.width))
        prepared = prepared.resize(
            (int(prepared.width * factor), int(prepared.height * factor)),
        )
    grayscale = ImageOps.grayscale(prepared)
    return ImageOps.autocontrast(grayscale)


def _run_tesseract(image: Any, language: str) -> tuple[str, str]:
    status = tesseract_status()
    if not status.get("available"):
        raise RuntimeError(
            "Tesseract OCR no está instalado o no fue detectado. "
            "Instale Tesseract OCR con el paquete de idioma español o defina TESSERACT_CMD."
        )
    import pytesseract

    command = str(status.get("command") or "")
    if command:
        pytesseract.pytesseract.tesseract_cmd = command
    prepared = _prepare_ocr_image(image)
    preferred_languages = [language]
    if "spa" not in preferred_languages:
        preferred_languages.append("spa")
    if "eng" not in preferred_languages:
        preferred_languages.append("eng")
    preferred_languages.append("")
    attempts = []
    for used_language in preferred_languages:
        attempts.append((used_language, "--oem 3 --psm 6"))
        attempts.append((used_language, "--oem 3 --psm 3"))
    last_error: Exception | None = None
    for used_language, config in attempts:
        try:
            kwargs: dict[str, Any] = {"config": config}
            if used_language:
                kwargs["lang"] = used_language
            value = pytesseract.image_to_string(prepared, **kwargs)
            return value, used_language or "predeterminado"
        except Exception as exc:
            last_error = exc
    raise RuntimeError(str(last_error or "No fue posible ejecutar Tesseract OCR"))


def _merge_digital_and_ocr(digital_text: str, ocr_text: str) -> str:
    digital = (digital_text or "").strip()
    ocr = (ocr_text or "").strip()
    if not digital:
        return ocr
    if not ocr:
        return digital
    digital_tokens = set(re.findall(r"[a-z0-9]{3,}", unicodedata.normalize("NFKD", digital.lower())))
    ocr_tokens = set(re.findall(r"[a-z0-9]{3,}", unicodedata.normalize("NFKD", ocr.lower())))
    novelty = len(ocr_tokens - digital_tokens) / max(1, len(ocr_tokens))
    if novelty < 0.08:
        return digital if len(digital) >= len(ocr) else ocr
    # OCR se coloca primero para conservar títulos, sellos y encabezados que suelen
    # perderse en la capa digital; el texto digital completa tablas y párrafos.
    return f"{ocr}\n\n[TEXTO DIGITAL COMPLEMENTARIO]\n{digital}"


def _extract_pdf(
    data: bytes,
    enable_ocr: bool = True,
    ocr_language: str = "spa+eng",
    max_ocr_pages: int = 50,
    ocr_mode: str = "all",
) -> tuple[str, dict[str, Any], list[str]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exc:
            raise ValueError("PDF protegido con contraseña") from exc

    digital_pages: list[str] = []
    for page in reader.pages:
        digital_pages.append(page.extract_text() or "")
    pages = list(digital_pages)
    warnings: list[str] = []
    ocr_pages: list[int] = []
    ocr_language_used = ""

    mode = str(ocr_mode or "all").strip().lower()
    if not enable_ocr:
        mode = "off"
    if mode not in {"all", "auto", "off"}:
        mode = "auto"

    if mode == "all":
        pages_needing_ocr = list(range(len(pages)))
    elif mode == "auto":
        pages_needing_ocr = [index for index, value in enumerate(pages) if len(value.strip()) < 160]
        # El encabezado de la primera página es decisivo para diferenciar, por ejemplo,
        # un acta de comité de una simple mención a una tabla de retención.
        if pages and 0 not in pages_needing_ocr:
            pages_needing_ocr.insert(0, 0)
    else:
        pages_needing_ocr = []

    if pages_needing_ocr:
        status = tesseract_status()
        if not status.get("available"):
            warnings.append(
                "OCR solicitado, pero Tesseract no está disponible. Se utilizó únicamente el texto digital. "
                "Instale Tesseract OCR con idioma español o configure TESSERACT_CMD."
            )
        else:
            try:
                import fitz
                from PIL import Image

                document = fitz.open(stream=data, filetype="pdf")
                limited_pages = pages_needing_ocr[: max(1, int(max_ocr_pages))]
                for index in limited_pages:
                    page = document.load_page(index)
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(2.6, 2.6), colorspace=fitz.csRGB, alpha=False)
                    image = Image.open(io.BytesIO(pixmap.tobytes("png")))
                    ocr_text, ocr_language_used = _run_tesseract(image, ocr_language)
                    if ocr_text.strip():
                        pages[index] = _merge_digital_and_ocr(digital_pages[index], ocr_text)
                        ocr_pages.append(index + 1)
                document.close()
                if len(pages_needing_ocr) > len(limited_pages):
                    warnings.append(
                        f"El OCR se limitó a {len(limited_pages)} páginas; quedaron "
                        f"{len(pages_needing_ocr) - len(limited_pages)} sin OCR."
                    )
            except Exception as exc:
                warnings.append(f"No fue posible ejecutar OCR en el PDF: {exc}")

    remaining_empty = sum(not value.strip() for value in pages)
    if reader.pages and remaining_empty == len(reader.pages):
        warnings.append("El PDF no contiene texto extraíble y el OCR no produjo resultados.")
    elif remaining_empty:
        warnings.append(f"{remaining_empty} página(s) siguen sin texto después de la extracción/OCR.")

    page_blocks = [f"=== PÁGINA {number} ===\n{value}" for number, value in enumerate(pages, start=1)]
    metadata = {
        "pages": len(reader.pages),
        "pdf_metadata": {str(k): str(v) for k, v in (reader.metadata or {}).items()},
        "ocr_requested_mode": mode,
        "ocr_performed": bool(ocr_pages),
        "ocr_pages": ocr_pages,
        "ocr_pages_requested": [index + 1 for index in pages_needing_ocr],
        "ocr_coverage_percent": round(len(ocr_pages) / max(1, len(reader.pages)) * 100, 1),
        "ocr_language": ocr_language_used,
        "digital_text_pages": sum(bool(value.strip()) for value in digital_pages),
    }
    return "\n\n".join(page_blocks), metadata, warnings



def _extract_docx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    table_cells = 0
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            table_cells += len(values)
            if any(values):
                parts.append(" | ".join(values))
    properties = document.core_properties
    return "\n".join(parts), {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "table_cells": table_cells,
        "created": properties.created.isoformat() if properties.created else "",
        "modified": properties.modified.isoformat() if properties.modified else "",
    }, []


def _extract_xlsx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    cells_read = 0
    truncated = False
    for ws in workbook.worksheets:
        parts.append(f"HOJA: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v not in (None, "")]
            cells_read += len(values)
            if values:
                parts.append(" | ".join(values))
            if cells_read >= 50_000:
                truncated = True
                break
        if truncated:
            break
    warnings = ["La hoja se truncó a 50.000 celdas para proteger el rendimiento."] if truncated else []
    return "\n".join(parts), {
        "sheets": workbook.sheetnames,
        "cells_read": cells_read,
        "created": workbook.properties.created.isoformat() if workbook.properties.created else "",
        "modified": workbook.properties.modified.isoformat() if workbook.properties.modified else "",
    }, warnings


def _extract_xls(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    import pandas as pd

    workbook = pd.ExcelFile(io.BytesIO(data))
    parts = []
    rows = 0
    for sheet in workbook.sheet_names:
        frame = workbook.parse(sheet_name=sheet, header=None, dtype=str).fillna("")
        parts.append(f"HOJA: {sheet}")
        for values in frame.itertuples(index=False, name=None):
            line = " | ".join(str(v).strip() for v in values if str(v).strip())
            if line:
                parts.append(line)
                rows += 1
            if rows >= 20_000:
                break
    warnings = ["La hoja se truncó a 20.000 filas."] if rows >= 20_000 else []
    return "\n".join(parts), {"sheets": workbook.sheet_names, "rows_read": rows}, warnings


def _extract_pptx(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts = []
    for number, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    slide_parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide_parts:
            parts.append(f"DIAPOSITIVA {number}\n" + "\n".join(slide_parts))
    properties = presentation.core_properties
    return "\n\n".join(parts), {
        "slides": len(presentation.slides),
        "created": properties.created.isoformat() if properties.created else "",
        "modified": properties.modified.isoformat() if properties.modified else "",
    }, []


def _extract_odt(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    from bs4 import BeautifulSoup

    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        xml = archive.read("content.xml")
    soup = BeautifulSoup(xml, "xml")
    text = "\n".join(node.get_text(" ", strip=True) for node in soup.find_all(["text:p", "text:h"]))
    return text, {}, []


def _extract_email(data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    message = BytesParser(policy=policy.default).parsebytes(data)
    parts = [
        f"Asunto: {message.get('subject', '')}",
        f"De: {message.get('from', '')}",
        f"Para: {message.get('to', '')}",
        f"Fecha: {message.get('date', '')}",
    ]
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain" and part.get_content_disposition() != "attachment":
                try:
                    parts.append(part.get_content())
                except Exception:
                    payload = part.get_payload(decode=True) or b""
                    parts.append(_decode_text(payload))
    else:
        try:
            parts.append(message.get_content())
        except Exception:
            parts.append(_decode_text(message.get_payload(decode=True) or b""))
    return "\n".join(str(x) for x in parts), {"subject": str(message.get("subject", ""))}, []


def _extract_image_ocr(data: bytes, ocr_language: str = "spa+eng") -> tuple[str, dict[str, Any], list[str]]:
    from PIL import Image

    image = Image.open(io.BytesIO(data))
    exif = image.getexif()
    exif_created = ""
    for tag in (36867, 306):
        if exif.get(tag):
            exif_created = str(exif.get(tag))
            break
    try:
        text, language_used = _run_tesseract(image, ocr_language)
        warning = [] if text.strip() else ["OCR ejecutado, pero no se detectó texto."]
        return text, {
            "width": image.width,
            "height": image.height,
            "ocr": True,
            "ocr_language": language_used,
            "created": exif_created,
        }, warning
    except Exception as exc:
        return "", {"width": image.width, "height": image.height, "ocr": False, "created": exif_created}, [
            f"OCR no disponible ({exc}). Instale Tesseract con el idioma español."
        ]


def _extract_plain(ext: str, data: bytes) -> tuple[str, dict[str, Any], list[str]]:
    text = _decode_text(data)
    if ext in {".html", ".htm", ".xml"}:
        try:
            from bs4 import BeautifulSoup

            text = BeautifulSoup(text, "html.parser").get_text("\n", strip=True)
        except Exception:
            pass
    elif ext == ".json":
        try:
            text = json.dumps(json.loads(text), ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            pass
    elif ext == ".rtf":
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
        text = text.replace("{", " ").replace("}", " ")
    return text, {}, []


def extract_document(
    payload: DocumentPayload,
    enable_ocr: bool = True,
    ocr_language: str = "spa+eng",
    max_ocr_pages: int = 50,
    ocr_mode: str = "all",
) -> ExtractedDocument:
    ext = payload.extension
    warnings: list[str] = []
    metadata: dict[str, Any] = {
        "bytes": len(payload.data),
        "source_modified_at": payload.source_modified_at,
    }
    try:
        if ext == ".pdf":
            text, extra, warnings = _extract_pdf(
                payload.data,
                enable_ocr=enable_ocr,
                ocr_language=ocr_language,
                max_ocr_pages=max_ocr_pages,
                ocr_mode=ocr_mode,
            )
        elif ext == ".docx":
            text, extra, warnings = _extract_docx(payload.data)
        elif ext in {".xlsx", ".xlsm"}:
            text, extra, warnings = _extract_xlsx(payload.data)
        elif ext == ".xls":
            text, extra, warnings = _extract_xls(payload.data)
        elif ext == ".pptx":
            text, extra, warnings = _extract_pptx(payload.data)
        elif ext == ".odt":
            text, extra, warnings = _extract_odt(payload.data)
        elif ext == ".eml":
            text, extra, warnings = _extract_email(payload.data)
        elif ext in IMAGE_EXTENSIONS:
            if enable_ocr and str(ocr_mode or "all").lower() != "off":
                text, extra, warnings = _extract_image_ocr(payload.data, ocr_language=ocr_language)
            else:
                from PIL import Image

                image = Image.open(io.BytesIO(payload.data))
                text, extra, warnings = "", {"width": image.width, "height": image.height, "ocr": False}, [
                    "La imagen requiere OCR para extraer texto; el OCR estaba desactivado."
                ]
        elif ext in TEXT_EXTENSIONS:
            text, extra, warnings = _extract_plain(ext, payload.data)
        else:
            raise ValueError("tipo de archivo no compatible")
        metadata.update(extra)
    except Exception as exc:
        text = ""
        warnings.append(f"No se pudo extraer el contenido: {exc}")

    text = re.sub(r"[ \t]+", " ", text or "")
    text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]
        warnings.append(f"El texto fue truncado a {MAX_EXTRACTED_CHARS:,} caracteres.")
    if not text and not warnings:
        warnings.append("No se encontró texto utilizable.")

    return ExtractedDocument(
        name=payload.name,
        text=text,
        sha256=payload.sha256,
        extension=ext,
        metadata=metadata,
        warnings=warnings,
    )


def detect_document_year(
    filename: str,
    text: str,
    metadata: dict[str, Any] | None = None,
    default_year: int | str | None = None,
) -> dict[str, Any]:
    """Detecta el año por metadatos originales o por la fecha OCR más reciente."""

    current_year = datetime.now().year
    minimum_year = 1990
    maximum_year = current_year + 1
    pattern = re.compile(r"(?<!\d)(?:19|20)\d{2}(?!\d)")
    scores: Counter[int] = Counter()
    origins: dict[int, set[str]] = {}

    def add_candidates(value: Any, weight: int, origin: str) -> None:
        raw = str(value or "")
        matches = pattern.findall(raw)
        matches += re.findall(r"D:((?:19|20)\d{2})", raw)
        for match in matches:
            year = int(match)
            if minimum_year <= year <= maximum_year:
                scores[year] += weight
                origins.setdefault(year, set()).add(origin)

    add_candidates(filename, 22, "año en el nombre del archivo")

    metadata = metadata or {}

    def walk_metadata(value: Any, parent_key: str = "") -> None:
        if isinstance(value, dict):
            for key, child in value.items():
                walk_metadata(child, str(key).lower())
            return
        if isinstance(value, (list, tuple)):
            for child in value:
                walk_metadata(child, parent_key)
            return
        normalized_key = re.sub(r"[^a-z]", "", parent_key)
        if any(term in normalized_key for term in ["creationdate", "created", "datetimeoriginal"]):
            add_candidates(value, 20, "fecha de creación interna del archivo")
        elif "sourcemodifiedat" in normalized_key:
            add_candidates(value, 18, "fecha original registrada en el ZIP/navegador")
        elif "modified" in normalized_key or "moddate" in normalized_key:
            add_candidates(value, 8, "fecha de modificación interna")

    walk_metadata(metadata)

    normalized_text = "".join(
        character
        for character in unicodedata.normalize("NFKD", (text or "").lower())
        if not unicodedata.combining(character)
    )
    full_dates: list[tuple[int, int, int]] = []

    def add_full_date(year: int, month: int, day: int) -> None:
        try:
            parsed = datetime(year, month, day)
        except ValueError:
            return
        if minimum_year <= parsed.year <= maximum_year:
            full_dates.append((parsed.year, parsed.month, parsed.day))

    for day, month, year in re.findall(r"(?<!\d)(\d{1,2})[/-](\d{1,2})[/-]((?:19|20)\d{2})(?!\d)", normalized_text):
        add_full_date(int(year), int(month), int(day))
    for year, month, day in re.findall(r"(?<!\d)((?:19|20)\d{2})[/-](\d{1,2})[/-](\d{1,2})(?!\d)", normalized_text):
        add_full_date(int(year), int(month), int(day))

    months = {
        "enero": 1,
        "febrero": 2,
        "marzo": 3,
        "abril": 4,
        "mayo": 5,
        "junio": 6,
        "julio": 7,
        "agosto": 8,
        "septiembre": 9,
        "setiembre": 9,
        "octubre": 10,
        "noviembre": 11,
        "diciembre": 12,
    }
    month_pattern = "|".join(months)
    for day, month_name, year in re.findall(
        rf"(?<!\d)(\d{{1,2}})\s+de\s+({month_pattern})\s+de\s+((?:19|20)\d{{2}})(?!\d)",
        normalized_text,
    ):
        add_full_date(int(year), months[month_name], int(day))

    if full_dates:
        closest_date = max(full_dates)
        scores[closest_date[0]] += 24
        origins.setdefault(closest_date[0], set()).add(
            f"fecha más reciente detectada por OCR/contenido ({closest_date[2]:02d}/{closest_date[1]:02d}/{closest_date[0]})"
        )

    add_candidates(normalized_text[:3_500], 6, "año detectado por OCR/contenido inicial")
    add_candidates(normalized_text[3_500:25_000], 2, "año detectado en el contenido")

    if scores:
        best_year = max(scores, key=lambda year: (scores[year], year))
        total = sum(scores.values())
        strong_origins = " ".join(origins[best_year])
        bonus = 0.2 if any(term in strong_origins for term in ["nombre", "creación", "fecha más reciente"]) else 0.0
        confidence = min(1.0, scores[best_year] / max(total, 1) + bonus)
        return {
            "year": str(best_year),
            "source": ", ".join(sorted(origins[best_year])),
            "confidence": round(confidence, 3),
            "candidates": dict(sorted(scores.items(), key=lambda item: (-item[1], -item[0]))[:5]),
        }

    fallback = str(default_year or current_year)
    return {
        "year": fallback,
        "source": "año predeterminado (sin fecha detectada)",
        "confidence": 0.0,
        "candidates": {},
    }
